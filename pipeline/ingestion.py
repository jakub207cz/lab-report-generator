from __future__ import annotations

import io
import json
import mimetypes
import re
import unicodedata
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Set, Tuple, Union

import pandas as pd
import pypdf
from PIL import Image, ImageDraw, ImageFont
from docx import Document
from pptx import Presentation

from pipeline.schemas import DocumentChunk, FigureData, NormalizedIngestionResult, TableData

try:
    from openpyxl import load_workbook
    from openpyxl.utils.cell import get_column_letter, range_boundaries
except Exception:  # pragma: no cover
    load_workbook = None
    get_column_letter = None
    range_boundaries = None

try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except Exception:  # pragma: no cover
    pdfminer_extract_text = None

try:
    import pytesseract
except Exception:  # pragma: no cover
    pytesseract = None

try:
    import pypdfium2 as pdfium
except Exception:  # pragma: no cover
    pdfium = None


UNIT_PATTERN = re.compile(r"(?:\(([^)]+)\)|\[([^\]]+)\])")

ASSIGNMENT_SECTION_ALIASES: List[tuple[str, str, List[str]]] = [
    ("assignment", "Zadání", ["zadani", "zadání"]),
    ("theory", "Teoretický úvod", ["teoreticky uvod", "teorie"]),
    ("schema", "Schéma zapojení", ["schema zapojeni", "obrazek schema zapojeni", "schéma zapojení"]),
    ("procedure", "Postup měření", ["postup mereni", "pracovni postup", "postup měření"]),
    ("tables_example", "Příklad tabulek", ["priklad tabulek", "tabulky", "namerene a vypocitane hodnoty"]),
    ("calculation_example", "Příklad výpočtů", ["priklad vypoctu", "vypocty"]),
    (
        "expected_graphs",
        "Předpokládaný průběh grafů",
        [
            "predpokladany prubeh grafu",
            "predpokladany prubeh merenych charakteristik",
            "predpokladany prubeh merenych charakteristik",
            "graficke prubehy",
            "prubeh grafu",
            "prubeh merenych charakteristik",
        ],
    ),
    ("conclusion", "Závěr", ["zaver", "závěr"]),
]


@dataclass
class BinarySource:
    filename: str
    data: bytes
    section_hint: Optional[str] = None
    mime_type: Optional[str] = None


class IngestionPipeline:
    def __init__(self, enable_ocr: bool = True) -> None:
        self.enable_ocr = enable_ocr
        self._figure_counter = 1

    def ingest_sources(self, sources: Iterable[BinarySource]) -> NormalizedIngestionResult:
        source_list = list(sources)
        result = NormalizedIngestionResult()
        aggregated_metadata: Dict[str, Any] = {}

        for source in source_list:
            ext = Path(source.filename).suffix.lower()
            parser = self._resolve_parser(ext)
            if not parser:
                continue

            parsed = parser(source)
            result.chunks.extend(parsed.chunks)
            result.tables.extend(parsed.tables)
            result.figures.extend(parsed.figures)
            for meta_key, meta_value in parsed.metadata.items():
                if meta_key not in aggregated_metadata:
                    aggregated_metadata[meta_key] = meta_value
                    continue

                current_value = aggregated_metadata[meta_key]
                if isinstance(current_value, dict) and isinstance(meta_value, dict):
                    current_value.update(meta_value)
                elif isinstance(current_value, list) and isinstance(meta_value, list):
                    current_value.extend(meta_value)
                else:
                    aggregated_metadata[meta_key] = meta_value

        result.metadata = {
            "source_count": len(source_list),
            "chunk_count": len(result.chunks),
            "table_count": len(result.tables),
            "figure_count": len(result.figures),
            "ocr_enabled": self.enable_ocr,
        }
        result.metadata.update(aggregated_metadata)

        debug_payload = {
            "metadata": result.metadata,
            "chunks": [chunk.model_dump() for chunk in result.chunks],
            "tables": [table.model_dump() for table in result.tables],
            "figures": [figure.model_dump() for figure in result.figures],
        }
        print("[IngestionPipeline][DEBUG] Extracted payload:")
        print(json.dumps(debug_payload, ensure_ascii=False, indent=2))

        return result

    def ingest_streamlit_files(self, uploaded_files: Iterable[Any], section_hint: Optional[str] = None) -> NormalizedIngestionResult:
        sources: List[BinarySource] = []
        for up in uploaded_files or []:
            mime_type = getattr(up, "type", None)
            sources.append(
                BinarySource(
                    filename=up.name,
                    data=up.getvalue(),
                    section_hint=section_hint,
                    mime_type=mime_type,
                )
            )
        return self.ingest_sources(sources)

    def ingest_file_paths(self, file_paths: Iterable[Union[str, Path]], section_hint: Optional[str] = None) -> NormalizedIngestionResult:
        sources: List[BinarySource] = []
        for raw_path in file_paths or []:
            path = Path(raw_path)
            if not path.exists() or not path.is_file():
                continue
            mime_type, _ = mimetypes.guess_type(str(path))
            sources.append(
                BinarySource(
                    filename=path.name,
                    data=path.read_bytes(),
                    section_hint=section_hint,
                    mime_type=mime_type,
                )
            )
        return self.ingest_sources(sources)

    def ingest_file_path(self, file_path: Union[str, Path], section_hint: Optional[str] = None) -> NormalizedIngestionResult:
        return self.ingest_file_paths([file_path], section_hint=section_hint)

    def _resolve_parser(self, ext: str):
        if ext == ".docx":
            return self.docx_parser
        if ext in {".xlsx", ".xls"}:
            return self.xlsx_parser
        if ext in {".txt", ".csv"}:
            return self.text_parser
        if ext == ".pdf":
            return self.pdf_parser
        if ext == ".pptx":
            return self.pptx_parser
        if ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}:
            return self.image_handler
        return None

    def _next_figure_id(self) -> str:
        figure_id = f"FIG-{self._figure_counter:03d}"
        self._figure_counter += 1
        return figure_id

    def _units_from_headers(self, headers: List[str]) -> List[str]:
        units: List[str] = []
        for h in headers:
            match = UNIT_PATTERN.search(h or "")
            if not match:
                units.append("")
                continue
            units.append((match.group(1) or match.group(2) or "").strip())
        return units

    def _normalize_text(self, value: str) -> str:
        base = unicodedata.normalize("NFKD", value or "")
        no_accents = "".join(ch for ch in base if not unicodedata.combining(ch))
        lowered = no_accents.lower().strip()
        lowered = re.sub(r"\s+", " ", lowered)
        return lowered

    def _extract_assignment_sections(self, full_text: str) -> Dict[str, Dict[str, str]]:
        lines = [line.strip() for line in (full_text or "").splitlines() if line and line.strip()]
        if not lines:
            return {}

        normalized_aliases: List[tuple[str, str, List[str]]] = [
            (key, title, [self._normalize_text(alias) for alias in aliases])
            for key, title, aliases in ASSIGNMENT_SECTION_ALIASES
        ]

        markers: List[tuple[int, str, str]] = []
        for idx, line in enumerate(lines):
            normalized_line = self._normalize_text(line)
            for key, title, aliases in normalized_aliases:
                if any(
                    normalized_line == alias
                    or normalized_line.startswith(f"{alias}:")
                    or normalized_line.startswith(f"{alias} -")
                    or normalized_line.startswith(f"{alias} ")
                    for alias in aliases
                ):
                    if markers and markers[-1][1] == key:
                        continue
                    markers.append((idx, key, title))
                    break

        if not markers:
            return {}

        sections: Dict[str, Dict[str, str]] = {}
        for marker_idx, (start_idx, key, title) in enumerate(markers):
            end_idx = markers[marker_idx + 1][0] if marker_idx + 1 < len(markers) else len(lines)
            section_lines = lines[start_idx + 1 : end_idx]
            section_text = "\n".join(section_lines).strip()
            sections[key] = {
                "title": title,
                "text": section_text,
            }

        return sections

    def _detect_assignment_section_for_page(self, page_text: str) -> Optional[str]:
        normalized_page_text = self._normalize_text(page_text)
        if not normalized_page_text:
            return None

        for key, _, aliases in ASSIGNMENT_SECTION_ALIASES:
            normalized_aliases = [self._normalize_text(alias) for alias in aliases]
            if any(
                alias in normalized_page_text
                for alias in normalized_aliases
            ):
                return key

        return None

    def _extract_pdf_page_images(self, source: BinarySource, reader: pypdf.PdfReader) -> tuple[List[FigureData], List[Dict[str, Any]]]:
        figures: List[FigureData] = []
        saved_figures: List[Dict[str, Any]] = []

        file_stem = Path(source.filename).stem
        output_dir = Path.cwd() / "extracted_assets" / file_stem
        output_dir.mkdir(parents=True, exist_ok=True)

        for page_idx, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "")
            section_key = self._detect_assignment_section_for_page(page_text) if (source.section_hint or "").lower() == "assignment" else None
            section_hint = f"assignment:{section_key}" if section_key else source.section_hint

            try:
                page_images = list(getattr(page, "images", []) or [])
            except Exception:
                page_images = []

            for image_idx, image_obj in enumerate(page_images, start=1):
                try:
                    image_name = getattr(image_obj, "name", f"page-{page_idx:03d}-img-{image_idx:03d}.png")
                    image_data = getattr(image_obj, "data", None)
                    if not image_data:
                        continue

                    image_suffix = Path(image_name).suffix or ".png"
                    save_name = f"page-{page_idx:03d}-img-{image_idx:03d}{image_suffix}"
                    save_path = output_dir / save_name
                    save_path.write_bytes(image_data)

                    figure_id = self._next_figure_id()
                    figures.append(
                        FigureData(
                            figure_id=figure_id,
                            source_file=source.filename,
                            page=page_idx,
                            section_hint=section_hint,
                            ocr_text=None,
                            confidence=1.0,
                        )
                    )
                    saved_figures.append(
                        {
                            "figure_id": figure_id,
                            "page": page_idx,
                            "path": str(save_path),
                        }
                    )
                except Exception:
                    continue

        return figures, saved_figures

    def _is_cell_empty(self, value: Any) -> bool:
        if value is None:
            return True
        if isinstance(value, str) and not value.strip():
            return True
        return False

    def _extract_excel_table_blocks(self, ws) -> List[Tuple[int, int, int, int]]:
        """
        Najde samostatné bloky neprázdných buněk v worksheetu.
        Blok = souvislá komponenta buněk (8-směrné sousedství).
        Vrací tuple: (min_row, max_row, min_col, max_col).
        """
        non_empty: Set[Tuple[int, int]] = set()
        for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
            for cell in row:
                if not self._is_cell_empty(cell.value):
                    non_empty.add((cell.row, cell.column))

        if not non_empty:
            return []

        visited: Set[Tuple[int, int]] = set()
        blocks: List[Tuple[int, int, int, int]] = []

        for start in non_empty:
            if start in visited:
                continue

            stack = [start]
            visited.add(start)
            component: List[Tuple[int, int]] = []

            while stack:
                r, c = stack.pop()
                component.append((r, c))

                for dr in (-1, 0, 1):
                    for dc in (-1, 0, 1):
                        if dr == 0 and dc == 0:
                            continue
                        nb = (r + dr, c + dc)
                        if nb in non_empty and nb not in visited:
                            visited.add(nb)
                            stack.append(nb)

            rows = [rc[0] for rc in component]
            cols = [rc[1] for rc in component]
            min_row, max_row = min(rows), max(rows)
            min_col, max_col = min(cols), max(cols)

            # Odfiltruj drobné artefakty (osamělé buňky)
            area_non_empty = len(component)
            if area_non_empty < 2:
                continue

            blocks.append((min_row, max_row, min_col, max_col))

        # Stabilní pořadí: shora dolů, zleva doprava
        blocks.sort(key=lambda b: (b[0], b[2]))
        return blocks

    def _sheet_name_with_range(self, sheet_name: str, min_row: int, max_row: int, min_col: int, max_col: int) -> str:
        if get_column_letter is None:
            return sheet_name
        return f"{sheet_name}!{get_column_letter(min_col)}{min_row}:{get_column_letter(max_col)}{max_row}"

    def _parse_excel_ref(self, ref_formula: str) -> Optional[Tuple[str, int, int, int, int]]:
        if not ref_formula or range_boundaries is None:
            return None
        if "!" not in ref_formula:
            return None

        sheet_part, range_part = ref_formula.split("!", 1)
        sheet_name = sheet_part.strip().strip("'").replace("''", "'")
        first_area = range_part.split(",", 1)[0].replace("$", "")

        try:
            min_col, min_row, max_col, max_row = range_boundaries(first_area)
        except Exception:
            return None
        return sheet_name, min_col, min_row, max_col, max_row

    def _extract_values_from_ref(self, workbook, ref_formula: Optional[str]) -> List[Any]:
        if not ref_formula:
            return []

        parsed = self._parse_excel_ref(ref_formula)
        if not parsed:
            return []

        sheet_name, min_col, min_row, max_col, max_row = parsed
        if sheet_name not in workbook.sheetnames:
            return []

        ws = workbook[sheet_name]
        values: List[Any] = []
        for row in range(min_row, max_row + 1):
            for col in range(min_col, max_col + 1):
                values.append(ws.cell(row=row, column=col).value)
        return values

    def _as_float(self, value: Any) -> Optional[float]:
        if value is None:
            return None
        if isinstance(value, (int, float)):
            return float(value)

        raw = str(value).strip().replace(" ", "")
        if not raw:
            return None
        if re.fullmatch(r"[+-]?\d+(?:[\.,]\d+)?(?:[eE][+-]?\d+)?", raw):
            try:
                return float(raw.replace(",", "."))
            except Exception:
                return None
        return None

    def _extract_chart_text(self, text_obj: Any) -> str:
        if text_obj is None:
            return ""
        if isinstance(text_obj, str):
            return text_obj.strip()

        tx = getattr(text_obj, "tx", None)
        if tx is not None:
            rich = getattr(tx, "rich", None)
            if rich is not None:
                parts: List[str] = []
                for paragraph in getattr(rich, "p", []) or []:
                    for run in getattr(paragraph, "r", []) or []:
                        t = getattr(run, "t", "")
                        if t:
                            parts.append(str(t))
                joined = "".join(parts).strip()
                if joined:
                    return joined

        value = getattr(text_obj, "v", None)
        if value is not None:
            return str(value).strip()
        return ""

    def _extract_series_name(self, workbook, series: Any, idx: int) -> str:
        title_obj = getattr(series, "title", None)
        if title_obj is not None:
            ref_formula = getattr(getattr(title_obj, "strRef", None), "f", None)
            if ref_formula:
                vals = self._extract_values_from_ref(workbook, ref_formula)
                if vals and vals[0] is not None:
                    return str(vals[0]).strip()

            title_value = getattr(title_obj, "v", None)
            if title_value is not None:
                return str(title_value).strip()

            extracted = self._extract_chart_text(title_obj)
            if extracted:
                return extracted

        return f"Série {idx}"

    def _extract_series_xy(self, workbook, series: Any) -> Tuple[List[float], List[float]]:
        y_ref = (
            getattr(getattr(getattr(series, "yVal", None), "numRef", None), "f", None)
            or getattr(getattr(getattr(series, "val", None), "numRef", None), "f", None)
            or getattr(getattr(getattr(series, "val", None), "strRef", None), "f", None)
        )
        x_ref = (
            getattr(getattr(getattr(series, "xVal", None), "numRef", None), "f", None)
            or getattr(getattr(getattr(series, "cat", None), "numRef", None), "f", None)
            or getattr(getattr(getattr(series, "cat", None), "strRef", None), "f", None)
        )

        y_raw = self._extract_values_from_ref(workbook, y_ref)
        x_raw = self._extract_values_from_ref(workbook, x_ref)
        if not y_raw:
            return [], []

        x_numeric: List[float] = []
        y_numeric: List[float] = []
        fallback_x = list(range(1, len(y_raw) + 1))

        for idx, y_val in enumerate(y_raw):
            y_num = self._as_float(y_val)
            if y_num is None:
                continue

            x_source = x_raw[idx] if idx < len(x_raw) else fallback_x[idx]
            x_num = self._as_float(x_source)
            if x_num is None:
                x_num = float(fallback_x[idx])

            x_numeric.append(x_num)
            y_numeric.append(y_num)

        return x_numeric, y_numeric

    def _render_chart_png(self, chart: Any, chart_label: str, workbook, save_path: Path) -> bool:
        series_data: List[Dict[str, Any]] = []
        for idx, series in enumerate(getattr(chart, "ser", []) or [], start=1):
            x_vals, y_vals = self._extract_series_xy(workbook, series)
            if not y_vals:
                continue
            series_data.append(
                {
                    "name": self._extract_series_name(workbook, series, idx),
                    "x": x_vals,
                    "y": y_vals,
                }
            )

        if not series_data:
            return False

        width, height = 1400, 900
        margin_left, margin_right, margin_top, margin_bottom = 130, 70, 150, 130
        plot_w = width - margin_left - margin_right
        plot_h = height - margin_top - margin_bottom

        xs = [x for s in series_data for x in s["x"]]
        ys = [y for s in series_data for y in s["y"]]
        if not xs or not ys:
            return False

        x_min, x_max = min(xs), max(xs)
        y_min, y_max = min(ys), max(ys)
        if x_min == x_max:
            x_min -= 1.0
            x_max += 1.0
        if y_min == y_max:
            y_min -= 1.0
            y_max += 1.0

        image = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(image)

        def _load_font(size: int) -> ImageFont.ImageFont:
            for font_path in [
                "C:/Windows/Fonts/arial.ttf",
                "C:/Windows/Fonts/calibri.ttf",
                "C:/Windows/Fonts/segoeui.ttf",
            ]:
                try:
                    return ImageFont.truetype(font_path, size)
                except Exception:
                    continue
            try:
                return ImageFont.truetype("DejaVuSans.ttf", size)
            except Exception:
                return ImageFont.load_default()

        title_font = _load_font(38)
        axis_font = _load_font(24)
        tick_font = _load_font(18)
        legend_font = _load_font(20)

        def _text_size(text: str, font: ImageFont.ImageFont) -> Tuple[int, int]:
            bbox = draw.textbbox((0, 0), text, font=font)
            return max(0, bbox[2] - bbox[0]), max(0, bbox[3] - bbox[1])

        def _wrap_text_to_width(text: str, font: ImageFont.ImageFont, max_width: int, max_lines: int = 3) -> List[str]:
            words = (text or "").split()
            if not words:
                return [""]

            lines: List[str] = []
            current = words[0]
            for word in words[1:]:
                candidate = f"{current} {word}".strip()
                w, _ = _text_size(candidate, font)
                if w <= max_width:
                    current = candidate
                else:
                    lines.append(current)
                    current = word
            lines.append(current)

            if len(lines) > max_lines:
                lines = lines[:max_lines]
                last = lines[-1]
                while last and _text_size(last + "…", font)[0] > max_width:
                    last = last[:-1]
                lines[-1] = (last + "…") if last else "…"

            return lines

        x0, y0 = margin_left, height - margin_bottom
        x1, y1 = width - margin_right, margin_top
        draw.line((x0, y0, x1, y0), fill="black", width=2)
        draw.line((x0, y0, x0, y1), fill="black", width=2)

        x_ticks = 12
        y_ticks = 12
        x_minor_ticks = 24
        y_minor_ticks = 24

        # Hustá mřížka: nejdřív vedlejší (minor), pak hlavní (major).
        for i in range(0, y_minor_ticks + 1):
            gy = margin_top + i * (plot_h / y_minor_ticks)
            draw.line((margin_left, gy, width - margin_right, gy), fill="#efefef", width=1)

        for i in range(0, x_minor_ticks + 1):
            gx = margin_left + i * (plot_w / x_minor_ticks)
            draw.line((gx, margin_top, gx, height - margin_bottom), fill="#f1f1f1", width=1)

        for i in range(0, y_ticks + 1):
            gy = margin_top + i * (plot_h / y_ticks)
            draw.line((margin_left, gy, width - margin_right, gy), fill="#d4d4d4", width=1)

        for i in range(0, x_ticks + 1):
            gx = margin_left + i * (plot_w / x_ticks)
            draw.line((gx, margin_top, gx, height - margin_bottom), fill="#d8d8d8", width=1)

        def map_x(x: float) -> float:
            return margin_left + ((x - x_min) / (x_max - x_min)) * plot_w

        def map_y(y: float) -> float:
            return (height - margin_bottom) - ((y - y_min) / (y_max - y_min)) * plot_h

        for i in range(x_ticks + 1):
            ratio = i / x_ticks
            x_pos = margin_left + ratio * plot_w
            draw.line((x_pos, y0, x_pos, y0 + 8), fill="black", width=2)
            x_val = x_min + ratio * (x_max - x_min)
            label = f"{x_val:.3f}".rstrip("0").rstrip(".")
            tw, th = _text_size(label, tick_font)
            draw.text((x_pos - tw / 2, y0 + 12), label, fill="black", font=tick_font)

        for i in range(y_ticks + 1):
            ratio = i / y_ticks
            y_pos = y0 - ratio * plot_h
            draw.line((x0 - 8, y_pos, x0, y_pos), fill="black", width=2)
            y_val = y_min + ratio * (y_max - y_min)
            label = f"{y_val:.3f}".rstrip("0").rstrip(".")
            tw, th = _text_size(label, tick_font)
            draw.text((x0 - 14 - tw, y_pos - th / 2), label, fill="black", font=tick_font)

        colors = ["#1f77b4", "#d62728", "#2ca02c", "#9467bd", "#ff7f0e", "#17becf"]

        for s_idx, series in enumerate(series_data):
            ordered = sorted(zip(series["x"], series["y"]), key=lambda p: p[0])
            pts = [(map_x(x), map_y(y)) for x, y in ordered]
            if len(pts) >= 2:
                draw.line(pts, fill=colors[s_idx % len(colors)], width=3)
            for px, py in pts:
                size = 7
                c = colors[s_idx % len(colors)]
                # Marker jako kolmé + (90°), ne diagonální X.
                draw.line((px - size, py, px + size, py), fill=c, width=2)
                draw.line((px, py - size, px, py + size), fill=c, width=2)

        title = self._extract_chart_text(getattr(chart, "title", None)) or chart_label
        title_lines = _wrap_text_to_width(title, title_font, max_width=plot_w)
        line_h = _text_size("Ag", title_font)[1] + 6
        title_top = 24
        for li, line in enumerate(title_lines):
            line_w, _ = _text_size(line, title_font)
            draw.text(((width - line_w) / 2, title_top + li * line_h), line, fill="black", font=title_font)

        x_title = self._extract_chart_text(getattr(getattr(chart, "x_axis", None), "title", None))
        y_title = self._extract_chart_text(getattr(getattr(chart, "y_axis", None), "title", None))
        if x_title:
            tw, th = _text_size(x_title, axis_font)
            draw.text((margin_left + (plot_w - tw) / 2, height - 52), x_title, fill="black", font=axis_font)
        if y_title:
            # Vertikální popisek osy Y se středem na ose.
            ty_w, ty_h = _text_size(y_title, axis_font)
            y_img = Image.new("RGBA", (ty_w + 10, ty_h + 10), (255, 255, 255, 0))
            y_draw = ImageDraw.Draw(y_img)
            y_draw.text((5, 5), y_title, fill="black", font=axis_font)
            y_rot = y_img.rotate(90, expand=True)
            y_pos_x = 18
            y_pos_y = int(margin_top + (plot_h - y_rot.height) / 2)
            image.paste(y_rot, (y_pos_x, y_pos_y), y_rot)

        legend_x = width - margin_right - 260
        legend_y = margin_top + 16
        for s_idx, series in enumerate(series_data):
            c = colors[s_idx % len(colors)]
            yy = legend_y + s_idx * 30
            draw.rectangle((legend_x, yy + 4, legend_x + 16, yy + 14), fill=c)
            draw.text((legend_x + 24, yy), str(series["name"]), fill="black", font=legend_font)

        image.save(save_path, format="PNG")
        return save_path.exists()

    def _extract_xlsx_chart_images_with_openpyxl(self, source: BinarySource, workbook) -> tuple[List[FigureData], List[Dict[str, Any]]]:
        figures: List[FigureData] = []
        saved_figures: List[Dict[str, Any]] = []

        file_stem = Path(source.filename).stem
        output_dir = Path.cwd() / "extracted_assets" / file_stem
        output_dir.mkdir(parents=True, exist_ok=True)

        chart_items: List[Tuple[str, str, Any]] = []
        for chart_sheet in getattr(workbook, "chartsheets", []):
            for idx, chart in enumerate(getattr(chart_sheet, "_charts", []) or [], start=1):
                chart_items.append(("chartsheet", f"{chart_sheet.title} #{idx}", chart))

        for ws in workbook.worksheets:
            for idx, chart in enumerate(getattr(ws, "_charts", []) or [], start=1):
                chart_items.append(("worksheet", f"{ws.title} #{idx}", chart))

        for idx, (kind, label, chart) in enumerate(chart_items, start=1):
            save_name = f"xlsx-chart-openpyxl-{idx:03d}.png"
            save_path = output_dir / save_name

            rendered = False
            try:
                rendered = self._render_chart_png(chart=chart, chart_label=label, workbook=workbook, save_path=save_path)
            except Exception:
                rendered = False

            figure_id = self._next_figure_id()
            prefix = "Chartsheet" if kind == "chartsheet" else "Worksheet chart"
            figures.append(
                FigureData(
                    figure_id=figure_id,
                    source_file=source.filename,
                    section_hint=source.section_hint,
                    ocr_text=f"{prefix}: {label}",
                    confidence=1.0 if rendered else 0.7,
                )
            )

            if rendered:
                saved_figures.append(
                    {
                        "figure_id": figure_id,
                        "page": None,
                        "path": str(save_path),
                    }
                )

        return figures, saved_figures

    def docx_parser(self, source: BinarySource) -> NormalizedIngestionResult:
        doc = Document(io.BytesIO(source.data))
        out = NormalizedIngestionResult()

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = (para.style.name or "").lower() if para.style else ""
            chunk_type = "heading" if "heading" in style_name else "paragraph"
            out.chunks.append(
                DocumentChunk(
                    text=text,
                    type=chunk_type,
                    source_file=source.filename,
                    section_hint=source.section_hint,
                    confidence=1.0,
                )
            )

        for table in doc.tables:
            rows = [[(cell.text or "").strip() for cell in row.cells] for row in table.rows]
            if not rows:
                continue
            headers = rows[0]
            body = rows[1:] if len(rows) > 1 else []
            out.tables.append(
                TableData(
                    headers=headers,
                    rows=body,
                    units=self._units_from_headers(headers),
                    source_file=source.filename,
                    section_hint=source.section_hint,
                )
            )
            out.chunks.append(
                DocumentChunk(
                    text=f"Tabulka: {' | '.join(headers)}",
                    type="table",
                    source_file=source.filename,
                    section_hint=source.section_hint,
                    confidence=1.0,
                )
            )

        image_reltype_suffix = "/image"
        for rel in doc.part.rels.values():
            if not rel.reltype.endswith(image_reltype_suffix):
                continue
            figure_id = self._next_figure_id()
            ocr_text = None
            confidence = 1.0
            blob = rel.target_part.blob
            if self.enable_ocr and pytesseract is not None:
                try:
                    img = Image.open(io.BytesIO(blob))
                    ocr_text = (pytesseract.image_to_string(img, lang="ces+eng") or "").strip() or None
                except Exception:
                    confidence = 0.7

            out.figures.append(
                FigureData(
                    figure_id=figure_id,
                    source_file=source.filename,
                    section_hint=source.section_hint,
                    ocr_text=ocr_text,
                    confidence=confidence,
                )
            )
            out.chunks.append(
                DocumentChunk(
                    text=f"Figure {figure_id} z DOCX",
                    type="figure",
                    source_file=source.filename,
                    section_hint=source.section_hint,
                    confidence=confidence,
                )
            )

        return out

    def text_parser(self, source: BinarySource) -> NormalizedIngestionResult:
        out = NormalizedIngestionResult()
        decoded_text = ""
        for encoding in ("utf-8", "cp1250", "latin-1"):
            try:
                decoded_text = source.data.decode(encoding)
                break
            except Exception:
                continue

        decoded_text = decoded_text.strip()
        if not decoded_text:
            return out

        out.chunks.append(
            DocumentChunk(
                text=decoded_text,
                type="paragraph",
                source_file=source.filename,
                section_hint=source.section_hint,
                confidence=1.0,
            )
        )

        if source.filename.lower().endswith(".csv"):
            try:
                df = pd.read_csv(io.BytesIO(source.data))
                headers = [str(c) for c in df.columns]
                rows = df.fillna("").astype(object).values.tolist()
                out.tables.append(
                    TableData(
                        headers=headers,
                        rows=rows,
                        units=self._units_from_headers(headers),
                        source_file=source.filename,
                        section_hint=source.section_hint,
                    )
                )
            except Exception:
                pass

        return out

    def xlsx_parser(self, source: BinarySource) -> NormalizedIngestionResult:
        out = NormalizedIngestionResult()
        table_layouts: List[Dict[str, Any]] = []

        if load_workbook is not None:
            try:
                wb = load_workbook(io.BytesIO(source.data), data_only=True)

                for ws in wb.worksheets:
                    blocks = self._extract_excel_table_blocks(ws)
                    for table_idx, (min_row, max_row, min_col, max_col) in enumerate(blocks, start=1):
                        matrix: List[List[Any]] = []
                        for r in range(min_row, max_row + 1):
                            row_values: List[Any] = []
                            for c in range(min_col, max_col + 1):
                                val = ws.cell(row=r, column=c).value
                                row_values.append("" if val is None else val)
                            matrix.append(row_values)

                        if not matrix:
                            continue

                        headers = [str(v).strip() if v is not None else "" for v in matrix[0]]
                        rows = [[str(v) if v is not None else "" for v in row] for row in matrix[1:]]

                        sheet_with_range = self._sheet_name_with_range(ws.title, min_row, max_row, min_col, max_col)
                        out.tables.append(
                            TableData(
                                headers=headers,
                                rows=rows,
                                units=self._units_from_headers(headers),
                                source_file=source.filename,
                                sheet_name=sheet_with_range,
                                section_hint=source.section_hint,
                            )
                        )
                        out.chunks.append(
                            DocumentChunk(
                                text=(
                                    f"XLSX tabulka {table_idx} v listu '{ws.title}' "
                                    f"({max_row - min_row + 1} řádků, {max_col - min_col + 1} sloupců)."
                                ),
                                type="table",
                                source_file=source.filename,
                                section_hint=source.section_hint,
                                confidence=1.0,
                            )
                        )

                        table_layouts.append(
                            {
                                "sheet_name": ws.title,
                                "sheet_range": sheet_with_range,
                                "start_row": min_row,
                                "end_row": max_row,
                                "start_col": min_col,
                                "end_col": max_col,
                                "row_count": max_row - min_row + 1,
                                "col_count": max_col - min_col + 1,
                            }
                        )

                reconstructed_figures, reconstructed_saved = self._extract_xlsx_chart_images_with_openpyxl(source, wb)
                out.figures.extend(reconstructed_figures)
                for fig in reconstructed_figures:
                    out.chunks.append(
                        DocumentChunk(
                            text=f"Figure {fig.figure_id} rekonstruovaný z XLSX ({fig.ocr_text or 'chart'})",
                            type="figure",
                            source_file=source.filename,
                            section_hint=source.section_hint,
                            confidence=fig.confidence,
                        )
                    )
                if reconstructed_saved:
                    out.metadata.setdefault("saved_figures", [])
                    out.metadata["saved_figures"].extend(reconstructed_saved)

                for ws in wb.worksheets:

                    for idx, _img in enumerate(getattr(ws, "_images", []) or [], start=1):
                        figure_id = self._next_figure_id()
                        out.figures.append(
                            FigureData(
                                figure_id=figure_id,
                                source_file=source.filename,
                                section_hint=source.section_hint,
                                ocr_text=f"Worksheet image: {ws.title} #{idx}",
                                confidence=0.95,
                            )
                        )
                        out.chunks.append(
                            DocumentChunk(
                                text=f"Figure {figure_id} z XLSX worksheet '{ws.title}' (image {idx})",
                                type="figure",
                                source_file=source.filename,
                                section_hint=source.section_hint,
                                confidence=0.95,
                            )
                        )
            except Exception:
                pass

        if table_layouts:
            out.metadata.setdefault("table_layouts", [])
            out.metadata["table_layouts"].extend(table_layouts)

        return out

    def pdf_parser(self, source: BinarySource) -> NormalizedIngestionResult:
        out = NormalizedIngestionResult()
        extracted_any_text = False
        extracted_text_parts: List[str] = []

        try:
            reader = pypdf.PdfReader(io.BytesIO(source.data))
            for idx, page in enumerate(reader.pages, start=1):
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    extracted_any_text = True
                    extracted_text_parts.append(page_text)
                    out.chunks.append(
                        DocumentChunk(
                            text=page_text,
                            type="paragraph",
                            source_file=source.filename,
                            page=idx,
                            section_hint=source.section_hint,
                            confidence=0.95,
                        )
                    )

            extracted_figures, saved_figures = self._extract_pdf_page_images(source, reader)
            out.figures.extend(extracted_figures)
            if saved_figures:
                out.metadata.setdefault("saved_figures", [])
                out.metadata["saved_figures"].extend(saved_figures)
        except Exception:
            pass

        if not extracted_any_text and pdfminer_extract_text is not None:
            try:
                text = (pdfminer_extract_text(io.BytesIO(source.data)) or "").strip()
                if text:
                    extracted_any_text = True
                    extracted_text_parts.append(text)
                    out.chunks.append(
                        DocumentChunk(
                            text=text,
                            type="paragraph",
                            source_file=source.filename,
                            section_hint=source.section_hint,
                            confidence=0.9,
                        )
                    )
            except Exception:
                pass

        # OCR fallback: pouze pokud není textová vrstva nebo extrakce selhala
        if not extracted_any_text and self.enable_ocr:
            ocr_chunks, ocr_figures = self._pdf_ocr_fallback(source)
            out.chunks.extend(ocr_chunks)
            out.figures.extend(ocr_figures)

        is_assignment_pdf = (source.section_hint or "").lower() == "assignment"
        if extracted_any_text and is_assignment_pdf:
            sections = self._extract_assignment_sections("\n\n".join(extracted_text_parts))
            if sections:
                out.metadata["assignment_sections"] = sections
                for section_key, section_data in sections.items():
                    out.chunks.append(
                        DocumentChunk(
                            text=section_data["title"],
                            type="heading",
                            source_file=source.filename,
                            section_hint=f"assignment:{section_key}",
                            confidence=0.95,
                        )
                    )
                    if section_data["text"]:
                        out.chunks.append(
                            DocumentChunk(
                                text=section_data["text"],
                                type="paragraph",
                                source_file=source.filename,
                                section_hint=f"assignment:{section_key}",
                                confidence=0.95,
                            )
                        )

        return out

    def _pdf_ocr_fallback(self, source: BinarySource) -> tuple[List[DocumentChunk], List[FigureData]]:
        chunks: List[DocumentChunk] = []
        figures: List[FigureData] = []

        if pytesseract is None or pdfium is None:
            chunks.append(
                DocumentChunk(
                    text="OCR fallback není dostupný (chybí pytesseract nebo pypdfium2).",
                    type="paragraph",
                    source_file=source.filename,
                    section_hint=source.section_hint,
                    confidence=0.3,
                )
            )
            return chunks, figures

        try:
            pdf = pdfium.PdfDocument(io.BytesIO(source.data))
        except Exception:
            chunks.append(
                DocumentChunk(
                    text="OCR fallback selhal při otevření PDF.",
                    type="paragraph",
                    source_file=source.filename,
                    section_hint=source.section_hint,
                    confidence=0.2,
                )
            )
            return chunks, figures

        for i in range(len(pdf)):
            try:
                page = pdf[i]
                pil_image = page.render(scale=2).to_pil()
                text = (pytesseract.image_to_string(pil_image, lang="ces+eng") or "").strip()
                figure_id = self._next_figure_id()

                figures.append(
                    FigureData(
                        figure_id=figure_id,
                        source_file=source.filename,
                        page=i + 1,
                        section_hint=source.section_hint,
                        ocr_text=text or None,
                        confidence=0.75 if text else 0.4,
                    )
                )

                if text:
                    chunks.append(
                        DocumentChunk(
                            text=text,
                            type="paragraph",
                            source_file=source.filename,
                            page=i + 1,
                            section_hint=source.section_hint,
                            confidence=0.75,
                        )
                    )
            except Exception:
                chunks.append(
                    DocumentChunk(
                        text=f"OCR selhal na straně {i + 1}.",
                        type="paragraph",
                        source_file=source.filename,
                        page=i + 1,
                        section_hint=source.section_hint,
                        confidence=0.2,
                    )
                )

        return chunks, figures

    def pptx_parser(self, source: BinarySource) -> NormalizedIngestionResult:
        out = NormalizedIngestionResult()
        prs = Presentation(io.BytesIO(source.data))

        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    text = (shape.text or "").strip()
                    if text:
                        out.chunks.append(
                            DocumentChunk(
                                text=text,
                                type="paragraph",
                                source_file=source.filename,
                                page=slide_idx,
                                section_hint=source.section_hint,
                                confidence=1.0,
                            )
                        )

                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    figure_id = self._next_figure_id()
                    ocr_text = None
                    conf = 1.0

                    if self.enable_ocr and pytesseract is not None:
                        try:
                            img = Image.open(io.BytesIO(shape.image.blob))
                            ocr_text = (pytesseract.image_to_string(img, lang="ces+eng") or "").strip() or None
                        except Exception:
                            conf = 0.7

                    out.figures.append(
                        FigureData(
                            figure_id=figure_id,
                            source_file=source.filename,
                            slide=slide_idx,
                            section_hint=source.section_hint,
                            ocr_text=ocr_text,
                            confidence=conf,
                        )
                    )
                    out.chunks.append(
                        DocumentChunk(
                            text=f"Figure {figure_id} ze slide {slide_idx}",
                            type="figure",
                            source_file=source.filename,
                            page=slide_idx,
                            section_hint=source.section_hint,
                            confidence=conf,
                        )
                    )

                if getattr(shape, "has_table", False) and shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([(cell.text or "").strip() for cell in row.cells])
                    if rows:
                        headers = rows[0]
                        body = rows[1:] if len(rows) > 1 else []
                        out.tables.append(
                            TableData(
                                headers=headers,
                                rows=body,
                                units=self._units_from_headers(headers),
                                source_file=source.filename,
                                page=slide_idx,
                                section_hint=source.section_hint,
                            )
                        )

        return out

    def image_handler(self, source: BinarySource) -> NormalizedIngestionResult:
        out = NormalizedIngestionResult()
        figure_id = self._next_figure_id()

        ocr_text = None
        conf = 1.0
        if self.enable_ocr and pytesseract is not None:
            try:
                image = Image.open(io.BytesIO(source.data))
                ocr_text = (pytesseract.image_to_string(image, lang="ces+eng") or "").strip() or None
            except Exception:
                conf = 0.6

        out.figures.append(
            FigureData(
                figure_id=figure_id,
                source_file=source.filename,
                section_hint=source.section_hint,
                ocr_text=ocr_text,
                confidence=conf,
            )
        )
        out.chunks.append(
            DocumentChunk(
                text=f"Figure {figure_id} ({source.filename})",
                type="figure",
                source_file=source.filename,
                section_hint=source.section_hint,
                confidence=conf,
            )
        )

        if ocr_text:
            out.chunks.append(
                DocumentChunk(
                    text=ocr_text,
                    type="paragraph",
                    source_file=source.filename,
                    section_hint=source.section_hint,
                    confidence=0.8,
                )
            )

        return out
